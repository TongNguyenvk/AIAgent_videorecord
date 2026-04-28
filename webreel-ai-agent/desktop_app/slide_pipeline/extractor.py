"""
Slide Extractor - Extract text, notes, and images from PPTX/PDF files.

Uses python-pptx for content extraction and LibreOffice for rendering.
Handles Docker environments with font packages and profile isolation.
"""

import json
import logging
import os
import subprocess
import tempfile
import uuid
from dataclasses import dataclass, field, asdict
from pathlib import Path
from typing import Optional

logger = logging.getLogger("slide_extractor")


@dataclass
class SlideData:
    """Data extracted from a single slide."""
    slide_number: int
    texts: list[str] = field(default_factory=list)
    notes: str = ""
    image_path: Optional[str] = None

    def to_dict(self) -> dict:
        return asdict(self)


def extract_text_from_pptx(pptx_path: Path) -> list[SlideData]:
    """Extract text and notes from each slide using python-pptx."""
    from pptx import Presentation

    prs = Presentation(str(pptx_path))
    slides_data = []

    for idx, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        texts.append(text)

        # Extract speaker notes
        notes = ""
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()

        slides_data.append(SlideData(
            slide_number=idx,
            texts=texts,
            notes=notes,
        ))

    logger.info(f"Extracted text from {len(slides_data)} slides")
    return slides_data


def render_slides_to_images(
    pptx_path: Path,
    output_dir: Path,
    job_id: str = None,
    target_width: int = 1920,
    target_height: int = 1080,
) -> list[Path]:
    """Render slides to PNG images using LibreOffice headless.

    Uses isolated user profile per job to prevent concurrent conflicts (Gotcha #2).
    """
    output_dir.mkdir(parents=True, exist_ok=True)

    # Unique profile dir to prevent LibreOffice lock conflicts
    profile_id = job_id or str(uuid.uuid4())[:8]
    profile_dir = Path(tempfile.mkdtemp(prefix=f"lo_profile_{profile_id}_"))

    try:
        # Convert PPTX to PDF first (more reliable for image extraction)
        cmd = [
            "soffice",
            f"-env:UserInstallation=file://{profile_dir}",
            "--headless",
            "--convert-to", "pdf",
            "--outdir", str(output_dir),
            str(pptx_path),
        ]

        logger.info(f"LibreOffice converting PPTX to PDF...")
        result = subprocess.run(
            cmd,
            capture_output=True, text=True, timeout=120,
        )

        if result.returncode != 0:
            logger.error(f"LibreOffice failed: {result.stderr[:500]}")
            raise RuntimeError(f"LibreOffice conversion failed: {result.stderr[:200]}")

        # Find the generated PDF
        pdf_name = pptx_path.stem + ".pdf"
        pdf_path = output_dir / pdf_name

        if not pdf_path.exists():
            raise FileNotFoundError(f"LibreOffice did not produce {pdf_path}")

        # Convert PDF pages to PNG images using pdf2image
        return _pdf_to_images(pdf_path, output_dir, target_width, target_height)

    finally:
        # Cleanup profile dir
        import shutil
        shutil.rmtree(profile_dir, ignore_errors=True)


def _pdf_to_images(
    pdf_path: Path,
    output_dir: Path,
    target_width: int = 1920,
    target_height: int = 1080,
) -> list[Path]:
    """Convert PDF pages to PNG images, normalized to target resolution (Gotcha #4)."""
    try:
        from pdf2image import convert_from_path

        images = convert_from_path(
            str(pdf_path),
            dpi=200,
            fmt="png",
        )
    except ImportError:
        logger.warning("pdf2image not available, trying ffmpeg fallback")
        return _pdf_to_images_ffmpeg(pdf_path, output_dir, target_width, target_height)

    image_paths = []
    for i, img in enumerate(images, 1):
        # Resize + pad to exact target resolution (Gotcha #4)
        img_resized = _resize_and_pad(img, target_width, target_height)
        img_path = output_dir / f"slide_{i:03d}.png"
        img_resized.save(str(img_path), "PNG")
        image_paths.append(img_path)
        logger.debug(f"Slide {i} -> {img_path}")

    logger.info(f"Rendered {len(image_paths)} slide images at {target_width}x{target_height}")
    return image_paths


def _resize_and_pad(img, target_w: int, target_h: int):
    """Resize image to fit within target dimensions, pad with black bars.

    Handles both 16:9 and 4:3 slides correctly (Gotcha #4).
    """
    from PIL import Image

    orig_w, orig_h = img.size
    ratio = min(target_w / orig_w, target_h / orig_h)
    new_w = int(orig_w * ratio)
    new_h = int(orig_h * ratio)

    # Resize maintaining aspect ratio
    img_resized = img.resize((new_w, new_h), Image.LANCZOS)

    # Create black canvas at target size
    canvas = Image.new("RGB", (target_w, target_h), (0, 0, 0))

    # Center the slide on canvas
    offset_x = (target_w - new_w) // 2
    offset_y = (target_h - new_h) // 2
    canvas.paste(img_resized, (offset_x, offset_y))

    return canvas


def _pdf_to_images_ffmpeg(
    pdf_path: Path,
    output_dir: Path,
    target_width: int,
    target_height: int,
) -> list[Path]:
    """Fallback: use ffmpeg to convert PDF to images."""
    cmd = [
        "ffmpeg", "-y",
        "-i", str(pdf_path),
        "-vf", f"scale={target_width}:{target_height}:force_original_aspect_ratio=decrease,"
               f"pad={target_width}:{target_height}:(ow-iw)/2:(oh-ih)/2:black",
        str(output_dir / "slide_%03d.png"),
    ]

    result = subprocess.run(cmd, capture_output=True, text=True, timeout=120)
    if result.returncode != 0:
        raise RuntimeError(f"ffmpeg PDF conversion failed: {result.stderr[:200]}")

    image_paths = sorted(output_dir.glob("slide_*.png"))
    logger.info(f"Rendered {len(image_paths)} slide images via ffmpeg")
    return image_paths


def extract_slides(
    file_path: Path,
    output_dir: Path,
    job_id: str = None,
) -> list[SlideData]:
    """Main entry point: extract text + render images from PPTX or PDF.

    Returns list of SlideData with text, notes, and image paths populated.
    """
    output_dir.mkdir(parents=True, exist_ok=True)
    images_dir = output_dir / "slides"
    images_dir.mkdir(parents=True, exist_ok=True)

    suffix = file_path.suffix.lower()

    if suffix in (".pptx", ".ppt"):
        # .ppt (binary) needs conversion to .pptx first for python-pptx
        actual_pptx = file_path
        pptx_tmp_dir = None
        if suffix == ".ppt":
            logger.info("Converting .ppt to .pptx via LibreOffice...")
            pptx_tmp_dir = Path(tempfile.mkdtemp(prefix="ppt_convert_"))
            profile_dir = Path(tempfile.mkdtemp(prefix="lo_ppt_"))
            try:
                cmd = [
                    "soffice",
                    f"-env:UserInstallation=file://{profile_dir}",
                    "--headless",
                    "--convert-to", "pptx",
                    "--outdir", str(pptx_tmp_dir),
                    str(file_path),
                ]
                result = subprocess.run(cmd, capture_output=True, text=True, timeout=120)
                if result.returncode != 0:
                    logger.warning(f"PPT->PPTX conversion failed: {result.stderr[:200]}")
                converted = pptx_tmp_dir / (file_path.stem + ".pptx")
                if converted.exists():
                    actual_pptx = converted
                    logger.info(f"Converted: {actual_pptx}")
                else:
                    logger.warning("PPT->PPTX conversion produced no output, skipping text extraction")
                    actual_pptx = None
            finally:
                import shutil
                shutil.rmtree(profile_dir, ignore_errors=True)

        # Extract text from PPTX (skip if conversion failed)
        if actual_pptx and actual_pptx.exists():
            try:
                slides_data = extract_text_from_pptx(actual_pptx)
            except Exception as e:
                logger.warning(f"Text extraction failed: {e}, continuing with images only")
                slides_data = None
        else:
            slides_data = None

        # Render images (from original .ppt or .pptx, LibreOffice handles both)
        image_paths = render_slides_to_images(file_path, images_dir, job_id=job_id)

        # Build slide data
        if slides_data and len(slides_data) == len(image_paths):
            for slide, img_path in zip(slides_data, image_paths):
                slide.image_path = str(img_path)
        else:
            # No text data or count mismatch, create from images
            slides_data = [
                SlideData(slide_number=i + 1, texts=[], notes="", image_path=str(p))
                for i, p in enumerate(image_paths)
            ]

        # Cleanup temp pptx
        if pptx_tmp_dir:
            import shutil
            shutil.rmtree(pptx_tmp_dir, ignore_errors=True)

        return slides_data

    elif suffix == ".pdf":
        # PDF: extract images only (no text extraction for now)
        image_paths = _pdf_to_images(
            file_path, images_dir,
            target_width=1920, target_height=1080,
        )

        return [
            SlideData(
                slide_number=i + 1,
                texts=[],
                notes="",
                image_path=str(img_path),
            )
            for i, img_path in enumerate(image_paths)
        ]

    else:
        raise ValueError(f"Unsupported file format: {suffix}. Use .pptx or .pdf")
